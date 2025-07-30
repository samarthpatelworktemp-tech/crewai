from crewai.tools import tool
from typing import Type, Literal
from pydantic import BaseModel, Field
from pptx import Presentation 

@tool
def generate_report(topic: str, summary: str, tasks: str, output_format: Literal["pptx"] = "pptx") -> str:
    """
    Generate a final report document (PDF, Word, or PowerPoint)
    using summary and task content. Returns the file path.
    """
 
    def _generate_ppt(self, topic, summary, tasks):
        ppt = Presentation()
 
        # Title Slide
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        slide.shapes.title.text = f"Project Report: {topic}"
 
        self._add_slide(ppt, "Summary", summary)
        self._add_slide(ppt, "Task Plan", tasks)
 
        file_path = f"{topic}_report.pptx"
        ppt.save(file_path)
        return f"✅ PPT report generated: {file_path}"
 
    def _add_slide(self, ppt, title, content):
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = title
        textbox = slide.placeholders[1]
        textbox.text = content.strip()[:3000]

    topic_clean = topic.replace(" ", "_")[:50]
    if output_format == "pptx":
        return _generate_ppt(topic_clean, summary, tasks)
    else:
        return "❌ Invalid format. Use 'pdf', 'docx', or 'pptx'."

    return [generate_report]