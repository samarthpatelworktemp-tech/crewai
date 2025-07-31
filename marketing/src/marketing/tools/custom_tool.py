from crewai.tools import BaseTool
from typing import Type, List
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches
import os


class PresentationToolInput(BaseModel):
    argument: str = Field(..., description=(
        "Input for the presentation tool. "
        "Expected format: 'topic|summary|target_audience|marketing_campaigns|campaigns|Budget_allocation|KPIs|Conclusion'"
    ))


class PresentationTool(BaseTool):
    name: str = "Presentation Creation Tool"
    description: str = (
        "A tool to create presentations based on marketing strategies and campaign ideas. "
        "Accepts structured string input and dynamically generates slides."
    )
    args_schema: Type[BaseModel] = PresentationToolInput

    def _run(self, argument: str) -> str:
        """
        Format:
        'topic|summary|target_audience|marketing_campaigns|campaigns|Budget_allocation|KPIs|Conclusion'
        """
        try:
            topic, summary, target_audience, marketing_campaigns, campaigns_str, budget_allocation, KPIs, conclusion = argument.split("|")
            campaigns = [c.strip() for c in campaigns_str.split(",") if c.strip()]
        except ValueError:
            return ("❌ Invalid argument format. Expected format: "
                    "'topic|summary|target_audience|marketing_campaigns|campaigns|Budget_allocation|KPIs|Conclusion'")

        return self._generate_ppt(
            topic.strip(), summary.strip(), target_audience.strip(),
            marketing_campaigns.strip(), campaigns, budget_allocation.strip(),
            KPIs.strip(), conclusion.strip()
        )

    def _generate_ppt(
        self, topic: str, summary: str, target_audience: str, marketing_campaigns: str,
        campaigns: List[str], budget_allocation: str, KPIs: str, conclusion: str
    ) -> str:

        image_folder = "charts/"
        image_extensions = [".png", ".jpg", ".jpeg", ".bmp", ".gif"]

        image_files = [f for f in os.listdir(image_folder) if os.path.splitext(f)[1].lower() in image_extensions]

        ppt = Presentation()

        # Title Slide
        title_slide_layout = ppt.slide_layouts[0]
        slide = ppt.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = f"Project Report: {topic}"

        # Content Slides
        self._add_slide(ppt, "Summary", summary)
        self._add_slide(ppt, "Target Audience", target_audience)
        self._add_slide(ppt, "Marketing Campaigns", marketing_campaigns)

        for i, campaign in enumerate(campaigns, start=1):
            self._add_slide(ppt, f"Campaign {i}", campaign)

        self._add_slide(ppt, "Budget Allocation", budget_allocation)
        self._add_slide(ppt, "KPIs", KPIs)
        self._add_slide(ppt, "Conclusion", conclusion)

        # Add image slides 
        for image_name in image_files:
            slide_layout = ppt.slide_layouts[1]
            slide = ppt.slides.add_slide(slide_layout)

            # Set slide title from image name
            slide.shapes.title.text = os.path.splitext(image_name)[0].replace('_', ' ').title()

            image_path = os.path.join(image_folder, image_name)
            pic = slide.shapes.add_picture(image_path, 0, 0)

            # Slide dimensions
            slide_width = ppt.slide_width
            slide_height = ppt.slide_height

            # Max display size
            max_width = Inches(6)
            max_height = Inches(4.5)

            # Scale image if needed
            if pic.width > max_width or pic.height > max_height:
                scale_x = max_width / pic.width
                scale_y = max_height / pic.height
                scale = min(scale_x, scale_y)
                pic.width = int(pic.width * scale)
                pic.height = int(pic.height * scale)

            # Center the image
            pic.left = int((slide_width - pic.width) / 2)
            pic.top = int((slide_height - pic.height) / 2)

        # Save presentation
        file_path = f"{topic.replace(' ', '_')}_report.pptx"
        ppt.save(file_path)

        return f"✅ PPT report generated: {file_path}"

    def _add_slide(self, ppt: Presentation, title: str, content: str):
        slide_layout = ppt.slide_layouts[1]  # Title and Content
        slide = ppt.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        textbox = slide.placeholders[1]
        textbox.text = content.strip()[:3000]
